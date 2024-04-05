"""
Created on Mon Jul 11 15:16:27 2022

@author: piotr.janczewski
"""

import os
from os.path import exists
from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE
# from pptx.dml.color import RGBColor
import pandas as pd
import datetime as dt
import streamlit as st


acn_login = os.getlogin()
folder = "C:/Users/" + acn_login + "/Accenture/Warsaw Analytics - Documents"
dest = "C:/Users/" + acn_login + "/Desktop/"

st.set_page_config(page_title="Warsaw AAI CV select",
                   # page_icon=folder + "/CVapp/images/favicon_accenture.png",
                   layout="wide",
                   initial_sidebar_state="expanded"
                   )

##############
# Input file locations and names
##############
CVpptxNm = "/01_CVs/Warsaw_Analytics.pptx"
CVpptxPt = folder + CVpptxNm
CVxlsxNm = "/01_CVs/Warsaw_Analytics CV refresh.xlsx"
CVxlsxPt = folder + CVxlsxNm
AVxlsxNm = "/Tracker/Tracker.xlsx"
AVxlsxPt = folder + AVxlsxNm

CVprs = Presentation(open(CVpptxPt, "rb"))
CVexl = pd.read_excel(CVxlsxPt, header=0)
AVexl = pd.read_excel(AVxlsxPt, 'AI Staffing Tracker', header=0)

Sel_txt = folder + '/01_CVs/CVapp/sel_list.txt'

##############
# Form parameters
##############

sen_5_val = False
sen_6_val = False
sen_7_val = False
sen_8_val = False
sen_9_val = False
sen_10_val = False
sen_11_val = False
sen_12_val = False

dpt_DS_val = False
dpt_DE_val = False
dpt_Oth_val = False

av_sl_val = 0


##############
# Scraping the pptx to produce a table with slidenums, names and positions
##############

CVlist = []
Sel_list = []

for slide in CVprs.slides:
    sldnm = CVprs.slides.index(slide)+1
    if slide.shapes.title != None:
        persn = slide.shapes.title.text.title()
        title = ''
        deptm = ''
        for shape in slide.shapes:
            if (shape.left > 5500000 and shape.left < 10000000 and
                shape.top < 500000):  # locating shape with job title
                if(len(shape.text.split(','))) > 1:
                    title = shape.text.split(',')[0]
                    deptm = shape.text.split(',')[1]
                else:
                    title = shape.text
                    deptm = ''
        if persn not in ['Team Cvs', 'Template']: 
            CVlist.append([sldnm, persn, title, deptm])

CVdf = pd.DataFrame(list(CVlist),
                    columns=['Slide', 'Person', 'Title', 'Dept'])


##############
# Functions
##############


def app_sec():
    #########
    # SIDEBAR #
    #########
    authorized_pple = ["tomasz.mostowski", "dominik.perykasza",
                       "piotr.janczewski", "aailead"]
    mot2pass = "JestesmyW0ln!"
    col1, mid, col2 = st.sidebar.columns([1, 4, 1])
    # mid.image(folder + "/CVapp/images/accenture_logo.png", width=200)
    username = st.sidebar.text_input("Login")
    pwd = st.sidebar.text_input("password", type='password')
    # username = 'aailead'; pwd ='JestesmyW0ln!' # For tests only
    if (username.lower() in authorized_pple) and (pwd == mot2pass):
        st.sidebar.success("Authorized")
        main_page()
    else:
        st.sidebar.error("Please enter valid credentials!")
        
def keepSlides(keepID, prs):

    # get slides to delete
    ids = [x for x in range(1, len(prs.slides._sldIdLst)+1) if x not in keepID]

    # subset report
    for i, slide in enumerate(prs.slides):
        # create slide dict
        id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}

        # iterate thorugh indexes
        if i+1 in ids:
            # get slide id
            slide_id = slide.slide_id

            # remove slide
            prs.part.drop_rel(id_dict[slide_id][1])
            del prs.slides._sldIdLst[id_dict[slide_id][0]]

    return prs


def kwdlookup(kwd_input, sld):
    kwds = kwd_input.lower()
    if 'and' in kwds:
        kwd_lgc = 'and'
        kwd_lst = kwds.split(" and")
    elif 'or' in kwd_input:
        kwd_lgc = 'or'
        kwd_lst = kwds.split(" or")
    else:
        kwd_lgc = 'na'
        kwd_lst = []
        kwd_lst.append(kwds)
    kwd_fnd = 0
    for i, kwd in enumerate(kwd_lst):
        for shape in CVprs.slides[sld].shapes:
            if hasattr(shape, "text") and (kwd in shape.text.lower()):
                kwd_fnd = kwd_fnd + 1
                break
    if kwd_lgc == 'and' and kwd_fnd == len(kwd_lst): kwd_fnd = 1
    elif kwd_lgc == 'or' and kwd_fnd > 0: kwd_fnd = 1
    elif kwd_lgc == 'na' and kwd_fnd == 1: kwd_fnd = 1
    elif kwds == '': kwd_fnd = 1

    return kwd_fnd


def main_page():

    global CVdf

    st.markdown("<p style='text-align:Left;font-family:Arial;" +
            "font-weight:bold;color:hsl(0, 100%, 0%); font-size:14px;'>" +
            "Specify criteria to export one-slider CVs</p>",
            unsafe_allow_html=True)

    st.write("")
    st.write("")

    with st.form("CV selection", clear_on_submit=False):

        av_sl_val = st.slider("Availability time (in weeks, counted" +
                              " since last Monday): 0=bench, 14=Over 3M",
                              min_value=0, max_value=14)
        st.write("")
        with st.expander("Seniority level"):
            sen_5_val = st.checkbox("5 - Associate Director")
            sen_6_val = st.checkbox("6 - Senior Manager")
            sen_7_val = st.checkbox("7 - Manager")
            sen_8_val = st.checkbox("8 - Associate Manager")
            sen_9_val = st.checkbox("9 - Team Lead/Consultant")
            sen_10_val = st.checkbox("10 - Senior Analyst")
            sen_11_val = st.checkbox("11 - Analyst")
            sen_12_val = st.checkbox("12 - Below Analyst")

        with st.expander("Department/profile"):
            dpt_DS_val = st.checkbox("Data Science")
            dpt_DE_val = st.checkbox("Data Engineering")
            dpt_Oth_val = st.checkbox("Other")
        
        with st.expander("Keywords"):
            st.markdown("<p style='font-size:12px;'>" +
                        "Use AND/OR (never both), eg.:<i>" +
                        " 'machine learning AND Azure AND risk advisory' or " +
                        " 'gcp OR Google Cloud Platform'</i></p>",
                        unsafe_allow_html=True)
            kwd_inp = st.text_input('Text to look up:')
            

        with st.expander("Full (preliminary) person list"):
            i = 0
            for Person in CVdf.Person:
                i = i + 1
                globals()['Person%s' % i] = st.checkbox(Person)

        
        # Premilinary listing before final export
        listed = st.form_submit_button("Filter people for final selection")
        if listed: 
            # Adding Level 
            lvl_dict = {
                "New Associate": 13,
                "Associate": 12,
                "Analyst": 11,
                "Senior Analyst": 10,
                "Consultant": 9,
                "Associate Manager": 8,
                "Manager": 7,
                "Digital Data Innovation Senior Principal": 7,
                "Senior Manager": 6,
                "Associate Director": 5,
                "Analytics Principal Director": 5,
                "Accenture leadership": 4
                }
            CVdf["Level"] = CVdf.Title.apply(lambda x: lvl_dict.get(x))
            
            # Adding availability

            CVdf["Avlbl"] = 14
            currweek = dt.date.today().isocalendar().week
            AVexl['AVweeks'] = AVexl['Availability week num']-currweek
            AVexl.loc[AVexl.AVweeks <= av_sl_val, 'AV'] = 1
            AVnmlst = list(AVexl['Full name'][AVexl.AV == 1])
            
            # Counting criteria selected
            
            sen_cnt = 0; dpt_cnt = 0; kwd_cnt = 0; prs_cnt = 0
            
            for sen_val in (sen_5_val, sen_6_val, sen_7_val, sen_8_val,
                            sen_9_val, sen_10_val, sen_11_val, sen_12_val):
                if sen_val: sen_cnt = sen_cnt + 1

            if (dpt_DS_val or dpt_DE_val or dpt_Oth_val): dpt_cnt = 1

            kwd_cnt = 0

            i = 0
            for Person in CVdf.Person:
                i = i + 1
                if globals()['Person%s' % i]: prs_cnt = prs_cnt + 1

            # Creating a name list for final selection

            CVdf['Select'] = 0
            for j in CVdf.Slide.index:
                slt = 0
                # Filter for Avlbl
                if CVdf.Person.values[j-1] in AVnmlst: slt = slt + 1

                # Filter for SenLvl
                i = 5; 
                for sen_val in (sen_5_val, sen_6_val, sen_7_val, sen_8_val,
                                sen_9_val, sen_10_val, sen_11_val, sen_12_val):
                    if (sen_val and CVdf.Level.values[j-1] == i): slt = slt + 1
                    i = i + 1
                if sen_cnt == 0: slt = slt + 1
                
                # Filter for Dept
                dptxt = CVdf.Dept.values[j-1][6:9].lower()
                if dpt_DS_val and dptxt == 'sci':  slt = slt + 1;
                if dpt_DE_val and dptxt == 'eng':  slt = slt + 1;
                if dpt_Oth_val and dptxt not in ['sci', 'eng']: slt = slt + 1
                if dpt_cnt == 0: slt = slt + 1
                
                # Filter for Keyword
                if kwdlookup(kwd_inp, j) == 1: slt = slt + 1

                # Filter for Person names
                if prs_cnt == 0: slt = slt + 1
                elif j > 0 and globals()['Person%s' % (j)] == True: slt = slt + 1

                CVdf.Select.values[j-1] = slt
                global Sel_list
                if slt == 5: Sel_list.append(CVdf.Person.values[j-1])
    
            with open(Sel_txt, 'w') as f:
                i = 0
                for line in Sel_list: 
                    i = i + 1; f.write(line)
                    if i < len(Sel_list): f.write('\n')

    # Displaying people for final approval
    if exists(Sel_txt):
        Sel_list = open(Sel_txt).read().split('\n')
        i = 0
        st.write("Filtered person list")
        for prsn in Sel_list:
            i = i + 1
            globals()['Select%s' % i] = st.checkbox(prsn + ' ', value = True)
            
    out_fn = st.text_input("Output file name (default: 'CVs free.pptx')")

    # Final export 
    exported = st.button("Export slides for the final people list")
    if exported: 
        i = 0; Fin_list = []
        for SelPrsn in Sel_list:
            i = i + 1
            if globals()['Select%s' % i] == True: Fin_list.append(SelPrsn)

        # Selecting final slides based on a name list
           
        KeepSLids = []
        for i, prsn in enumerate(CVdf.Person):
            if prsn in Fin_list: KeepSLids.append(CVdf.Slide[i])

        keepSlides(KeepSLids, CVprs)

        if out_fn == "": flnm = "CVs free.pptx"
        elif ".pptx" not in out_fn: flnm = out_fn + ".pptx"
        else: flnm = out_fn
        out_path = dest + flnm
        CVprs.save(out_path)      

        os.remove(Sel_txt)
            
    st.write("https://www.youtube.com/watch?v=WNnzw90vxrE of this form")

        
##############
# PAGE SET UP
##############


hide_streamlit_style = """
            <style>
            #MainMenu {visibility: hidden;}
            footer {visibility: hidden;}
            </style>
            """
st.markdown(hide_streamlit_style, unsafe_allow_html=True)

app_sec()