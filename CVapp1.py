"""
Created on Mon Jul 11 15:16:27 2022

@author: piotr.janczewski
"""

import os
import json
import openai
import pandas as pd
import datetime as dt
import streamlit as st
import unidecode as ud

from os.path import exists
from pptx import Presentation
from unidecode import unidecode

st.set_page_config(page_title="AI Industry Hub CV select",
                   # page_icon=folder + "/CVapp/images/favicon_accenture.png",
                   layout="wide",
                   initial_sidebar_state="expanded"
                   )

##############
# Input file locations and names
##############

acn_login = os.getlogin()
acn_path = "C:/Users/" + acn_login

CV_path = acn_path + "/Accenture/Warsaw Analytics - Documents/01_CVs/"
CV_flnm = "Warsaw_Analytics_FY23_template.pptx"
CV_file = CV_path + CV_flnm

AV_path = acn_path + "/Accenture/Staffing & productivity AI group - General/Dashboard/"
AV_flnm = "myScheduling_People_Extract.xlsx"
AV_file = AV_path + AV_flnm

LCR_flnm = "Staffing_data - LCR calc.xlsx"
LCR_file = AV_path + LCR_flnm

Promo_path = acn_path + "/Accenture/AI Executives Warsaw - General/"
Promo_flnm = "AI Ind Hub - promo slides.pptx"
Promo_file = Promo_path + Promo_flnm

App_path = acn_path + "/Desktop/genAI/CVapp"
Sel_txt = CV_path + '/CVapp/sel_list.txt'

dest_path = acn_path + "/Desktop/"

##############
# Dictionaries & initial parameters
##############

# Dictionary of shape names
shape_name_dict = {'Text Placeholder 1': 'About me',
                'Text Placeholder 2': 'Industry experience',
                'Text Placeholder 3': 'Name',
                'Text Placeholder 4': 'Level & role',
                'Text Placeholder 6': 'Education',
                'Text Placeholder 7': 'Key skills',
                'Text Placeholder 8': 'Languages',
                'Text Placeholder 9': 'Relevant project experience'} 

# Dictionary of titles per level
lvl_dict = {
    "13-New Associate": 13,
    "12-Associate": 12,
    "11-Analyst": 11,
    "10-Senior Analyst": 10,
    "9-Team Lead/Consultant": 9,
    "8-Associate Manager": 8,
    "7-Manager": 7,
    "7-Digital Data Innovation Senior Principal": 7,
    "6-Senior Manager": 6,
    "5-Associate Director": 5,
    "5-Analytics Principal Director": 5,
    "4-Accenture leadership": 4,
    "MD": 4
    }
sen_vars = {}
for i in range(4, 13):
    sen_vars[f'sen_{i}'] = False

dpt_DS = False
dpt_DE = False
dpt_Oth = False

av_sl = 0


##############
# Functions
##############

# Scraping the pptx to produce a table with slidenums, names and positions
def scrap_CVs(CVprs):

    CVprs = Presentation(open(CV_file, "rb"))
    shape_list = []

    for slide in CVprs.slides:

        sldnm = CVprs.slides.index(slide)+1
        sldnt = slide.notes_slide
        if sldnt.notes_text_frame:
            sldnt_text = sldnt.notes_text_frame.text

        for shape in slide.shapes:
            shape_properties = {
            'sld_nm': sldnm,
            'sld_nt': sldnt_text,
            'name': shape.name
            }

            # Check if the shape is a text shape and if it has text
            if shape.has_text_frame:
                text = ""
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
                shape_properties['text'] = text.strip()  # Strip to remove trailing newline

            shape_list.append(shape_properties)
        shapes_df = pd.DataFrame(shape_list)

        # Map shape names to sections using shape_name_dict
        shapes_df['section'] = 'Other'
        shapes_df.loc[shapes_df['name'].str.startswith('Picture'), 'section'] = 'Picture'
        shapes_df.loc[shapes_df['name'].str.startswith('Text Pl'), 'section'] = shapes_df['name'].apply(lambda x: shape_name_dict.get(x, 'else'))

        names_df = shapes_df[['sld_nm', 'sld_nt']].drop_duplicates()
        names_df = names_df.applymap(lambda x: x.lower() if isinstance(x, str) else x)

    return shapes_df, names_df

# Function to fill missing Worker and Resource Name
def fill_missing_values(row):
    if pd.isna(row['EID']) and pd.notna(row['Enterprise ID']):
        row['EID'] = row['Enterprise ID']
    if pd.isna(row['Worker']) and pd.notna(row['EID']):
        row['Worker'] = row['EID'].title().replace('.', ' ')
    if pd.isna(row['Resource Name']):
        if pd.notna(row['Worker']):
            row['Resource Name'] = ', '.join(row['Worker'].split()[::-1])
    return row

# Loading input files
def load_inputs(AV_file, LCR_file, names_df): 

    AV_df = pd.read_excel(AV_file, 'People_Extract', header=0)
    AV_df = AV_df[['Resource Name', 'Resource Email', 'First Availability Date']].dropna()
    AV_df['EID'] = AV_df['Resource Email'].str.split('@').str[0]
    AV_df['Resource Name'] = AV_df['Resource Name'].str.title()

    All_df = pd.read_excel(LCR_file, 'LCR', header=0)
    All_df = pd.merge(All_df, AV_df, left_on='Enterprise ID', right_on='EID', how='outer')
    All_df['EID'] = All_df['EID'].fillna(All_df['Enterprise ID'])
    All_df = pd.merge(All_df, names_df, left_on='EID', right_on='sld_nt', how='outer')
    
    # Apply the function to fill missing values
    All_df = All_df.apply(fill_missing_values, axis=1).dropna(subset=['EID', 'First Availability Date'], how='all')
    All_df = All_df.sort_values('Resource Name')
    All_df = All_df[['Worker', 'Resource Name', 'EID', 'sld_nm', 'Management Level', 'People Lead', 'LCR in $', 'First Availability Date']].sort_values('Resource Name')
    
    return All_df

def keepSlides(keepID, prs):

    # get slides to delete
    ids = [x for x in range(1, len(prs.slides._sldIdLst)+1) if x not in keepID]

    # subset report
    for i, slide in enumerate(prs.slides):
        # create slide dict
        id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}

        # iterate thorugh indexes
        if i+1 in ids:
            # get slide id
            slide_id = slide.slide_id

            # remove slide
            prs.part.drop_rel(id_dict[slide_id][1])
            del prs.slides._sldIdLst[id_dict[slide_id][0]]

    return prs

# def kwdlookup(kwd_input, sld):
#     kwds = kwd_input.lower()
#     if 'and' in kwds:
#         kwd_lgc = 'and'
#         kwd_lst = kwds.split(" and")
#     elif 'or' in kwd_input:
#         kwd_lgc = 'or'
#         kwd_lst = kwds.split(" or")
#     else:
#         kwd_lgc = 'na'
#         kwd_lst = []
#         kwd_lst.append(kwds)
#     kwd_fnd = 0
#     for i, kwd in enumerate(kwd_lst):
#         for shape in CVprs.slides[sld].shapes:
#             if hasattr(shape, "text") and (kwd in shape.text.lower()):
#                 kwd_fnd = kwd_fnd + 1
#                 break
#     if kwd_lgc == 'and' and kwd_fnd == len(kwd_lst): kwd_fnd = 1
#     elif kwd_lgc == 'or' and kwd_fnd > 0: kwd_fnd = 1
#     elif kwd_lgc == 'na' and kwd_fnd == 1: kwd_fnd = 1
#     elif kwds == '': kwd_fnd = 1

#     return kwd_fnd


def main_page(All_df,shapes_df):

    st.markdown("Specify criteria to export one-slider CVs")
    st.write("")
    
    with st.form("CV selection", clear_on_submit=False):

        av_sl = st.slider("Availability time (in weeks, counted" +
                              " since last Monday): 0=bench, 14=Over 3M",
                              min_value=0, max_value=14)
        st.write("")
        with st.expander("Seniority level"):
            sen_4 = st.checkbox("4 - Accenture Leadership")
            sen_5 = st.checkbox("5 - Associate Director")
            sen_6 = st.checkbox("6 - Senior Manager")
            sen_7 = st.checkbox("7 - Manager")
            sen_8 = st.checkbox("8 - Associate Manager")
            sen_9 = st.checkbox("9 - Team Lead/Consultant")
            sen_10 = st.checkbox("10 - Senior Analyst")
            sen_11 = st.checkbox("11 - Analyst")
            sen_12 = st.checkbox("12 - Below Analyst")

        with st.expander("Department/profile"):
            dpt_DS = st.checkbox("Data Science")
            dpt_DE = st.checkbox("Data Engineering")
            dpt_Oth = st.checkbox("Other")
        
        with st.expander("Keywords"):
            st.markdown("<p style='font-size:12px;'>" +
                        "Use AND/OR (never both), eg.:<i>" +
                        " 'machine learning AND Azure AND risk advisory' or " +
                        " 'gcp OR Google Cloud Platform'</i></p>",
                        unsafe_allow_html=True)
            kwd_inp = st.text_input('Text to look up:')
            

        with st.expander("Preliminary person list"):
            i = 0
            for Person in All_df['Worker']:
                i = i + 1
                globals()['Person%s' % i] = st.checkbox(str(i) + ' ' + Person)
        
        # Premilinary listing before final export
        
        listed = st.form_submit_button("Filter people for final selection")
        if listed: 
            print('Success!')

            # Adding Level 
            All_df["Level"] = All_df['Management Level'].apply(lambda x: lvl_dict.get(x))
            
            # Adding availability
            All_df["Avlbl"] = 14
            currweek = dt.date.today().isocalendar().week
            All_df['Availability week num'] =  All_df['First Availability Date'].dt.isocalendar().week
            All_df['AVweeks'] = All_df['Availability week num']-currweek
            All_df.loc[All_df['AVweeks'] <= av_sl, 'AV'] = 1
            if av_sl == 14: AVnmlst = list(All_df['Worker'])
            else: AVnmlst = list(All_df['Worker'][All_df['AV'] == 1])
            print(All_df[All_df['AV']==1][['Worker', 'First Availability Date', 'Availability week num', 'AV']])
            ## Tu SKOŃCZYŁEM, aktualizując kolejne przyciski i testując streamlitem
            
            # Counting criteria selected
            
            sen_cnt = 0; dpt_cnt = 0; kwd_cnt = 0; prs_cnt = 0
            
            for sen in (sen_4, sen_5, sen_6, sen_7, sen_8,
                            sen_9, sen_10, sen_11, sen_12):
                if sen: sen_cnt = sen_cnt + 1

            if (dpt_DS or dpt_DE or dpt_Oth): dpt_cnt = 1

            kwd_cnt = 0

            # Creating a name list for final selection

            # All_df['Select'] = 0
            # for j in All_df['sld_nm'].index:
            #     slt = 0
            #     # Filter for Avlbl
            #     if All_df.Person.values[j-1] in AVnmlst: slt = slt + 1

            #     # Filter for SenLvl
            #     i = 5; 
            #     for sen in (sen_5, sen_6, sen_7, sen_8,
            #                     sen_9, sen_10, sen_11, sen_12):
            #         if (sen and All_df.Level.values[j-1] == i): slt = slt + 1
            #         i = i + 1
            #     if sen_cnt == 0: slt = slt + 1
                
            #     # Filter for Dept
            #     dptxt = All_df.Dept.values[j-1][6:9].lower()
            #     if dpt_DS and dptxt == 'sci':  slt = slt + 1;
            #     if dpt_DE and dptxt == 'eng':  slt = slt + 1;
            #     if dpt_Oth and dptxt not in ['sci', 'eng']: slt = slt + 1
            #     if dpt_cnt == 0: slt = slt + 1
                
            #     # Filter for Keyword
            #     if kwdlookup(kwd_inp, j) == 1: slt = slt + 1

            #     # Filter for Person names
            #     if prs_cnt == 0: slt = slt + 1
            #     elif j > 0 and globals()['Person%s' % (j)] == True: slt = slt + 1

            #     All_df.Select.values[j-1] = slt
            #     global Sel_list
            #     if slt == 5: Sel_list.append(All_df.Person.values[j-1])
                
            # with open(Sel_txt, 'w') as f:
            #     i = 0
            #     for line in Sel_list: 
            #         i = i + 1; f.write(line)
            #         if i < len(Sel_list): f.write('\n')

           # Displaying people for final approval

    # Final export 
    
    with st.form("CV export", clear_on_submit=False):
        if exists(Sel_txt):
            Sel_list = open(Sel_txt).read().split('\n')
            i = 0
            st.write("Filtered person list")
            for prsn in Sel_list:
                i = i + 1
                globals()['Select%s' % i] = st.checkbox(prsn + ' ', value = True)
        out_fn = st.text_input("Output file name (default: 'CVs free.pptx')")
        exported = st.form_submit_button("Export slides for the final people list")
        if exported: 
            i = 0; Fin_list = []
            for SelPrsn in Sel_list:
                i = i + 1
                if globals()['Select%s' % i] == True: Fin_list.append(SelPrsn)
    
            # Selecting final slides based on a name list
               
            KeepSLids = []
            for i, prsn in enumerate(All_df.Person):
                if prsn in Fin_list: KeepSLids.append(All_df.Slide[i])
    
            keepSlides(KeepSLids, CVprs)
    
            if out_fn == "": flnm = "CVs free.pptx"
            elif ".pptx" not in out_fn: flnm = out_fn + ".pptx"
            else: flnm = out_fn
            out_path = dest + flnm
            CVprs.save(out_path)      
    
            # os.remove(Sel_txt)
            
    st.write("https://www.youtube.com/watch?v=WNnzw90vxrE of this form")

        
##############
# PAGE SET UP
##############

hide_streamlit_style = """
            <style>
            #MainMenu {visibility: hidden;}
            footer {visibility: hidden;}
            </style>
            """

shapes_df, names_df = scrap_CVs(CV_file)
All_df = load_inputs(AV_file, LCR_file, names_df)

main_page(All_df, shapes_df)