"""
Created on Mon Jul 11 15:16:27 2022

@author: piotr.janczewski
@author: justyna.zbiegien
"""

import os
import openai
import pandas as pd
import unidecode as ud
import datetime as dt
import streamlit as st
from pptx import Presentation

st.set_page_config(page_title="AI Industry Hub CV select",
                   layout="wide",
                   initial_sidebar_state="expanded"
                   )


##############
# Input file locations and names - for local use
##############

# acn_login = os.getlogin()
# acn_path = "C:/Users/" + acn_login

# CV_path = acn_path + "/Accenture/Warsaw Analytics - Documents/01_CVs/"
# CV_flnm = "Warsaw_Analytics_FY23_template.pptx"
# CV_file = CV_path + CV_flnm

# AV_path = acn_path + "/Accenture/Staffing & productivity AI group - General/Dashboard/"
# AV_flnm = "myScheduling_People_Extract.xlsx"
# AV_file = AV_path + AV_flnm

# LCR_flnm = "Staffing_data - LCR calc.xlsx"
# LCR_file = AV_path + LCR_flnm

##############
# Input file locations and names - for online POC
##############

acn_path = "source samples/" 

CV_path = acn_path
CV_flnm = "Warsaw_Analytics_FY23.pptx"
CV_file = CV_path + CV_flnm

AV_path = acn_path
AV_flnm = "myScheduling_People_Extract.xlsx"
AV_file = AV_path + AV_flnm

LCR_flnm = "Staffing_data - LCR calc.xlsx"
LCR_file = AV_path + LCR_flnm

##############
# Dictionaries & initial parameters
##############

photo_adj_slides = 4
init_promo_slides = 2
final_promo_slides = 3

if 'filtered_df' not in st.session_state:
    st.session_state['filtered_df'] = pd.DataFrame()

if 'first_filtered' not in st.session_state:
    st.session_state['first_filtered'] = pd.DataFrame()


# Dictionary of shape names
shape_name_dict = {'Text Placeholder 1': 'About me',
                   'Text Placeholder 2': 'Industry experience',
                   'Text Placeholder 3': 'Name',
                   'Text Placeholder 4': 'Level & role',
                   'Text Placeholder 6': 'Education',
                   'Text Placeholder 7': 'Key skills',
                   'Text Placeholder 8': 'Languages',
                   'Text Placeholder 9': 'Relevant project experience'}

# Dictionary of titles per level
lvl_dict = {
    "13-New Associate": 13,
    "12-Associate": 12,
    "11-Analyst": 11,
    "10-Senior Analyst": 10,
    "9-Team Lead/Consultant": 9,
    "8-Associate Manager": 8,
    "7-Manager": 7,
    "7-Digital Data Innovation Senior Principal": 7,
    "6-Senior Manager": 6,
    "5-Associate Director": 5,
    "5-Analytics Principal Director": 5,
    "4-Accenture leadership": 4,
    "MD": 4
}
sen_vars = {}
for i in range(4, 13):
    sen_vars[f'sen_{i}'] = False

dpt_DS = False
dpt_DE = False
dpt_Oth = False

av_sl = 0

##############
# Functions
##############

# Scraping the pptx to produce a table with slidenums, names and positions
def scrap_CVs(CV_file):
    CVprs = Presentation(open(CV_file, "rb"))
    shape_list = []

    # Define the range of slides to exclude
    slides_to_exclude = set(range(1, init_promo_slides+1))
    fin_excl = photo_adj_slides + final_promo_slides
    slides_to_exclude = slides_to_exclude.union(set(range(len(CVprs.slides) - fin_excl + 1, len(CVprs.slides) + 1)))
    for slide in CVprs.slides:
        sld_nm = CVprs.slides.index(slide) + 1
        if sld_nm in slides_to_exclude:
            continue  # Skip slides to exclude
        sldnt = slide.notes_slide
        if sldnt.notes_text_frame:
            sldnt_text = sldnt.notes_text_frame.text

        # Check if the shape is a text shape and if it has text
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join([paragraph.text for paragraph in shape.text_frame.paragraphs])
                if "4" in shape.name:
                    if "data scien" in text.lower():
                        dept = "Data Science"
                    elif "data engineer" in text.lower():
                        dept = "Data Engineering"
                    else:
                        dept = "Other"

                    shape_properties = {
                        'sld_nm': sld_nm,
                        'sld_nt': sldnt_text,
                        'name': shape.name,
                        'Dept': dept,
                        'text': text.strip()
                    }

                    shape_list.append(shape_properties)

    shapes_df = pd.DataFrame(shape_list)

    # Map shape names to sections using shape_name_dict
    shapes_df['section'] = 'Other'
    shapes_df.loc[shapes_df['name'].str.startswith('Picture'), 'section'] = 'Picture'
    shapes_df.loc[shapes_df['name'].str.startswith('Text Pl'), 'section'] = shapes_df['name'].apply(
        lambda x: shape_name_dict.get(x, 'else'))

    names_df = shapes_df[['sld_nm', 'sld_nt', 'Dept']].drop_duplicates()
    names_df = names_df.applymap(lambda x: x.lower() if isinstance(x, str) else x)

    return shapes_df, names_df, CVprs


# Function to fill missing Worker and Resource Name
def fill_missing_values(row):
    if pd.isna(row['EID']) and pd.notna(row['Enterprise ID']):
        row['EID'] = row['Enterprise ID']
    if pd.isna(row['Worker']) and pd.notna(row['EID']):
        row['Worker'] = row['EID'].title().replace('.', ' ')
    if pd.isna(row['Resource Name']):
        if pd.notna(row['Worker']):
            row['Resource Name'] = ', '.join(row['Worker'].split()[::-1])
    return row


# Loading input files
def load_inputs(AV_file, LCR_file, names_df):
    AV_df = pd.read_excel(AV_file, 'People_Extract', header=0)
    AV_df = AV_df[['Resource Name', 'Resource Email', 'First Availability Date']].dropna()
    AV_df['EID'] = AV_df['Resource Email'].str.split('@').str[0]
    AV_df['Resource Name'] = AV_df['Resource Name'].str.title()

    All_df = pd.read_excel(LCR_file, 'LCR', header=0)
    All_df = pd.merge(All_df, AV_df, left_on='Enterprise ID', right_on='EID', how='outer')
    All_df['EID'] = All_df['EID'].fillna(All_df['Enterprise ID'])
    All_df = pd.merge(All_df, names_df, left_on='EID', right_on='sld_nt', how='outer')

    # Apply the function to fill missing values
    All_df = All_df.apply(fill_missing_values, axis=1).dropna(subset=['EID', 'First Availability Date'], how='all')
    All_df = All_df.sort_values('Resource Name')
    All_df = All_df[['Worker', 'Resource Name', 'EID', 'sld_nm', 'Management Level', 'People Lead', 'LCR in $',
                     'First Availability Date', 'Dept']].sort_values('Resource Name')

    All_df = All_df.dropna(subset=['sld_nm'])
    return All_df

# def kwdlookup(kwd_input, sld):
#     kwds = kwd_input.lower()
#     if 'and' in kwds:
#         kwd_lgc = 'and'
#         kwd_lst = kwds.split(" and")
#     elif 'or' in kwd_input:
#         kwd_lgc = 'or'
#         kwd_lst = kwds.split(" or")
#     else:
#         kwd_lgc = 'na'
#         kwd_lst = []
#         kwd_lst.append(kwds)
#     kwd_fnd = 0
#     for i, kwd in enumerate(kwd_lst):
#         for shape in CVprs.slides[sld].shapes:
#             if hasattr(shape, "text") and (kwd in shape.text.lower()):
#                 kwd_fnd = kwd_fnd + 1
#                 break
#     if kwd_lgc == 'and' and kwd_fnd == len(kwd_lst): kwd_fnd = 1
#     elif kwd_lgc == 'or' and kwd_fnd > 0: kwd_fnd = 1
#     elif kwd_lgc == 'na' and kwd_fnd == 1: kwd_fnd = 1
#     elif kwds == '': kwd_fnd = 1

#     return kwd_fnd

def initial_selection(All_df, shapes_df):
    st.markdown("Specify criteria to export one-slider CVs")
    st.write("")

    # with st.form("CV selection", clear_on_submit=False):

    av_sl = st.slider("Availability time (in weeks, counted" +
                        " since last Monday): 0=bench, 14=Over 3M",
                        min_value=0, max_value=14)
    st.write("")
    seniority_checks = {}
    with st.expander("Seniority level"):
        seniority_levels = [
            "4 - Accenture Leadership", 
            "5 - Associate Director", 
            "6 - Senior Manager",
            "7 - Manager", 
            "8 - Associate Manager", 
            "9 - Team Lead/Consultant",
            "10 - Senior Analyst", 
            "11 - Analyst", 
            "12 - Below Analyst"
        ]

        for level in seniority_levels:
            level_id = int(level.split(' - ')[0])
            seniority_checks[level_id] = st.checkbox(level)

    with st.expander("Department/profile"):
        dpt_DS = st.checkbox("Data Science")
        dpt_DE = st.checkbox("Data Engineering")
        dpt_Oth = st.checkbox("Other")

    # with st.expander("Keywords"):
    #     st.markdown("<p style='font-size:12px;'>" +
    #                 "Use AND/OR (never both), eg.:<i>" +
    #                 " 'machine learning AND Azure AND risk advisory' or " +
    #                 " 'gcp OR Google Cloud Platform'</i></p>",
    #                 unsafe_allow_html=True)
    #     kwd_inp = st.text_input('Text to look up:')
        kwd_inp = '' # Dummy to remove after proper kwd search is defined

    All_df.reset_index(drop=True, inplace=True)

    person_checks = {}
    with st.expander("Preliminary person list"):
        for index, Person in enumerate(All_df['Worker'], start=0):
            person_checks[index] = st.checkbox(f"{index} {Person}")

    # Premilinary listing before final export

    All_df['Select'] = False
    All_df['AV'] = 0

    listed = st.button("Filter people for final selection")
    if listed:
        filter_people(seniority_checks, person_checks, kwd_inp, dpt_DS, dpt_DE, dpt_Oth, av_sl)

    # Displaying people for final approval
    filtered_df = All_df[(All_df['Select'] == True) & (All_df['AV'] ==1)]

    if not filtered_df.empty:
        st.session_state['filtered_df'] = filtered_df
        st.session_state['first_filtered'] = filtered_df
    
    with st.expander("Filtered people list"):
        if not st.session_state['first_filtered'].empty:
            for index, row in st.session_state['first_filtered'].iterrows():
                checked=st.checkbox(f"{row['Worker']} - {row['Dept']} - Level {row['Level']} - AV {row['AVweeks']}", value=True)
                st.session_state['filtered_df'].loc[index, 'Select'] = checked

    if not st.session_state['filtered_df'].empty:
        st.session_state['filtered_df'] = st.session_state['filtered_df'][st.session_state['filtered_df']['Select'] == True]

    return st.session_state.get('filtered_df', pd.DataFrame())

def filter_people(seniority_checks, person_checks, kwd_inp, dpt_DS, dpt_DE, dpt_Oth, av_sl):
    active_filters_count =0
    # Counting active filters
    if any(seniority_checks.values()):
        active_filters_count += 1
    active_filters_count += (dpt_DS or dpt_DE or dpt_Oth) + bool(kwd_inp)

    # Adding Level
    All_df["Level"] = All_df['Management Level'].apply(lambda x: lvl_dict.get(x))

    # Adding availability
    All_df["Avlbl"] = 14
    currweek = dt.date.today().isocalendar().week
    All_df['Availability week num'] = All_df['First Availability Date'].dt.isocalendar().week
    All_df['AVweeks'] = All_df['Availability week num'] - currweek
    All_df.loc[All_df['AVweeks'] <= av_sl, 'AV'] = 1

    # Flitering criteria
    for index, row in All_df.iterrows():
        matched_criteria_count = 0

        # Counting criteria selected
        if row['Level'] in [lvl for lvl, checked in seniority_checks.items() if checked]:
                matched_criteria_count += 1

        if ((dpt_DS and 'sci' in row['Dept'].lower()) or
                (dpt_DE and 'eng' in row['Dept'].lower()) or
                (dpt_Oth and 'sci' not in row['Dept'].lower() and 'eng' not in row['Dept'].lower())):
            matched_criteria_count += 1

        # kwd_imp should be here

        All_df.loc[index, 'Select'] = matched_criteria_count == active_filters_count

        # Checking initial people list
        if person_checks.get(index):
            All_df.loc[index, 'Select'] = True
            All_df.loc[index, 'AV'] = 1

def export_to_excel(df, filepath):
    df = df.rename(columns={
    'Resource Name': 'Name',
    'Management Level': 'Position Level'
    })

    # Konwersja i formatowanie daty
    if 'First Availability Date' in df.columns:
        df['First Availability Date'] = pd.to_datetime(df['First Availability Date']).dt.strftime('%d.%m.%Y')

    # Filtruj DataFrame, aby zawierał tylko potrzebne kolumny
    columns_to_export = ['Name', 'EID', 'People Lead', 'Position Level', 'First Availability Date', 'LCR in $']
    
    # Sprawdzanie, czy wszystkie wymagane kolumny są w DataFrame
    if all(column in df.columns for column in columns_to_export):
        fil_df = df[columns_to_export]
        print(fil_df)
        
        # Eksportowanie do pliku Excel
        with pd.ExcelWriter(filepath, engine='openpyxl', mode='w') as writer:
            fil_df.to_excel(writer, index=False)
        print(f"Data exported successfully to {filepath}")
    else:
        missing_columns = [column for column in columns_to_export if column not in df.columns]
        print(f"Missing columns in DataFrame: {missing_columns}")

def remove_unwanted_slides(presentation, keep_slides_ids):
    """
    Usuwa slajdy, które nie znajdują się w podanym zbiorze identyfikatorów.
    Args:
    presentation (Presentation): Obiekt prezentacji.
    keep_slides_ids (set): Zbiór identyfikatorów slajdów do zachowania.
    """
    keep_slides_ids = set(map(lambda x: int(float(x)), keep_slides_ids)) # Konwersja na zestaw liczb całkowitych
    
    # Sprawdzenie, czy jest 1 lub więcej init_promo_slides, budowa std_sld_ids
    if init_promo_slides>0:
        standard_slides_ids = {1}
        for i in range(1,init_promo_slides +1):
                keep_slides_ids.update({i})
                
    fin_excl = photo_adj_slides + final_promo_slides
    
    for i in range(len(CVprs.slides) - fin_excl + 1, len(CVprs.slides) - final_promo_slides + 1):
        print(i)
        standard_slides_ids.update({i})
    keep_slides_ids.update(standard_slides_ids)

    # Uzyskujemy dostęp do listy identyfikatorów slajdów
    slide_ids = presentation.slides._sldIdLst
    # Iterujemy w odwrotnej kolejności, aby indeksy nie były zakłócane po usunięciu
    for i in reversed(range(len(slide_ids))):
        # Usuwamy slajd, jeśli jego indeks nie znajduje się w zbiorze do zachowania
        if i+1 not in keep_slides_ids:
            del slide_ids[i]

    return presentation
    
def create_presentation(filtered_df, presentation, output_path):
    keep_slides_ids = set(filtered_df['sld_nm'].astype(str))
    remove_unwanted_slides(presentation, keep_slides_ids)
    presentation.save(output_path)

def final_export(filtered_df):
    # Final export 

    with st.form("CV export", clear_on_submit=False):
        dest = st.text_input("Enter the directory path to save the CV and LCR files:", "path/to/directory")
        out_fn_CVs = st.text_input("Output CV file name", "AI Ind Hub CVs for ___.pptx")
        out_fn_LCR = st.text_input("Output LCR file name", "AI Ind Hub CVs for ___.xlsx")
        submit_button = st.form_submit_button("Export all slides for the filtered people list")
        
        if submit_button:
            output_path_CVs = f"{dest}/{out_fn_CVs if out_fn_CVs.endswith('.pptx') else out_fn_CVs + '.pptx'}"
            output_path_LCR = f"{dest}/{out_fn_LCR if out_fn_LCR.endswith('.xlsx') else out_fn_LCR + '.xlsx'}"
            create_presentation(filtered_df, CVprs, output_path_CVs)
            export_to_excel(filtered_df, output_path_LCR)
            st.success(f"Exported successfully to {dest}")
        else:
            st.error("Please filter the data before exporting.")

    st.markdown("For help, visit [YouTube](https://www.youtube.com/watch?v=WNnzw90vxrE)")

##############
# PAGE SET UP
##############

hide_streamlit_style = """
            <style>
            #MainMenu {visibility: hidden;}
            footer {visibility: hidden;}
            </style>
            """

shapes_df, names_df, CVprs = scrap_CVs(CV_file)
All_df = load_inputs(AV_file, LCR_file, names_df)

filtered_df = initial_selection(All_df, shapes_df)

final_export(filtered_df)